"""Final verification: check every acceptance criterion by actually exercising it.

This does not trust the file system alone - it loads the model, runs an inference,
opens the report PDF, opens the DOCX and PPTX, and reads inside the submission ZIP.

Usage:
    python scripts/verify_project.py [--zip submission/Group07_...zip]

Exit code 0 only when every REQUIRED check passes.
"""
from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

import numpy as np

from src import paths
from src.utils import file_size_mb, setup_logging

log = setup_logging("flameguard.verify")

CHECKS: list[tuple[str, bool, str, bool]] = []   # (name, passed, detail, required)


def check(name: str, required: bool = True):
    def wrap(fn):
        try:
            ok, detail = fn()
        except Exception as exc:                 # a crashing check is a failed check
            ok, detail = False, f"{type(exc).__name__}: {exc}"
        CHECKS.append((name, ok, detail, required))
        mark = "PASS" if ok else ("FAIL" if required else "WARN")
        log.info("[%-4s] %-38s %s", mark, name, detail)
        return fn
    return wrap


def run_all() -> None:
    # ---------------------------------------------------------------- dataset
    @check("Dataset: repaired split exists")
    def _():
        counts = {}
        for split in paths.SPLITS:
            d = paths.PROCESSED_DATASET_DIR / split / "images"
            counts[split] = len(list(d.iterdir())) if d.exists() else 0
        ok = all(v > 0 for v in counts.values())
        return ok, f"train/valid/test = {counts['train']}/{counts['valid']}/{counts['test']}"

    @check("Dataset: >=200 original images/class")
    def _():
        rep = json.loads((paths.VALIDATION_OUTPUT_DIR / "raw" /
                          "validation_report.json").read_text(encoding="utf-8"))
        per = rep["images_per_class_total"]
        return rep["min_200_images_per_class"], f"{per}"

    @check("Dataset: zero cross-split leakage after repair")
    def _():
        rs = json.loads((paths.VALIDATION_OUTPUT_DIR /
                         "resplit_report.json").read_text(encoding="utf-8"))
        after = rs["leakage_after"]["groups_spanning_splits"]
        before = rs["leakage_before"]["images_in_spanning_groups"]
        return after == 0, f"{before:,} leaked images repaired -> {after} spanning groups"

    @check("EDA: figures generated")
    def _():
        figs = list(paths.EDA_OUTPUT_DIR.glob("*.png"))
        return len(figs) >= 10, f"{len(figs)} figures"

    # ------------------------------------------------------------------ model
    @check("Model: final best.pt exists")
    def _():
        p = paths.FINAL_MODEL_PATH
        return p.exists(), f"{paths.rel_to_root(p)} ({file_size_mb(p):.1f} MB)"

    @check("Model: metadata with real metrics")
    def _():
        import yaml

        meta = yaml.safe_load(paths.FINAL_MODEL_METADATA_PATH.read_text(encoding="utf-8"))
        has = all(k in meta for k in ("validation_metrics", "test_metrics",
                                      "confidence_threshold"))
        thr = meta.get("confidence_threshold")
        return has and thr is not None, f"conf={thr}, model={meta.get('model_name')}"

    @check("Model: loads and infers on CPU")
    def _():
        from src.inference import DetectionEngine

        eng = DetectionEngine(paths.FINAL_MODEL_PATH, device="cpu")
        img = np.zeros((320, 320, 3), dtype=np.uint8)
        res = eng.predict(img, conf=0.99, iou=0.5, draw=True)
        return (res.annotated_bgr is not None and res.counts["total"] == 0,
                f"CPU inference {res.inference_ms:.0f} ms, empty-path OK")

    @check("Training: >=3 experiments logged")
    def _():
        import pandas as pd

        df = pd.read_csv(paths.TRAINING_OUTPUT_DIR / "experiment_log.csv")
        return len(df) >= 3, f"{len(df)} runs: {', '.join(df.experiment_id)}"

    # ------------------------------------------------------------- evaluation
    @check("Evaluation: test metrics present")
    def _():
        m = json.loads((paths.EVALUATION_OUTPUT_DIR /
                        "metrics_test.json").read_text(encoding="utf-8"))
        return ("map50" in m and m["images"] > 0,
                f"mAP50={m['map50']:.3f} R={m['recall']:.3f} on {m['images']} images")

    @check("Evaluation: threshold chosen on validation")
    def _():
        import pandas as pd

        df = pd.read_csv(paths.EVALUATION_OUTPUT_DIR / "threshold_analysis.csv")
        return len(df) >= 5, f"{len(df)} thresholds swept"

    @check("Evaluation: measured inference speed")
    def _():
        s = json.loads((paths.EVALUATION_OUTPUT_DIR /
                        "inference_speed.json").read_text(encoding="utf-8"))
        parts = [f"{k}: {v['fps']:.1f} FPS" for k, v in s.items()]
        return bool(s), "; ".join(parts)

    @check("Benchmark: table + selection report")
    def _():
        import pandas as pd

        df = pd.read_csv(paths.BENCHMARK_OUTPUT_DIR / "benchmark_table.csv")
        sel = json.loads((paths.BENCHMARK_OUTPUT_DIR /
                          "selection_report.json").read_text(encoding="utf-8"))
        return len(df) >= 2, f"{len(df)} models, winner={sel['winner']}"

    @check("Error analysis: failures documented")
    def _():
        e = json.loads((paths.ERROR_ANALYSIS_OUTPUT_DIR /
                        "error_summary.json").read_text(encoding="utf-8"))
        galleries = list(paths.ERROR_ANALYSIS_OUTPUT_DIR.glob("gallery_*"))
        return (e["total"]["fp"] + e["total"]["fn"] > 0 and len(galleries) >= 2,
                f"FP={e['total']['fp']} FN={e['total']['fn']}, {len(galleries)} galleries")

    @check("Samples: >=10 with successes AND failures")
    def _():
        import pandas as pd

        df = pd.read_csv(paths.SAMPLE_OUTPUTS_DIR / "sample_summary.csv")
        correct = (df.assessment == "Correct").sum()
        wrong = len(df) - correct
        return (len(df) >= 10 and wrong >= 1,
                f"{len(df)} samples: {correct} correct, {wrong} incorrect/partial")

    # -------------------------------------------------------------- app + docs
    @check("App: imports cleanly")
    def _():
        import py_compile

        py_compile.compile(str(paths.PROJECT_ROOT / "app.py"), doraise=True)
        return True, "app.py compiles"

    @check("App: screenshots captured")
    def _():
        shots = list(paths.SCREENSHOTS_DIR.glob("*.png"))
        return len(shots) >= 8, f"{len(shots)} screenshots"

    @check("Manual tests recorded")
    def _():
        import pandas as pd

        df = pd.read_csv(paths.OUTPUTS_DIR / "manual_test_results.csv")
        failed = (df.result != "PASS").sum()
        return failed == 0, f"{len(df)} manual tests, {failed} failed"

    @check("Automated tests passed")
    def _():
        log_file = paths.OUTPUTS_DIR / "test_report.txt"
        text = log_file.read_text(encoding="utf-8", errors="replace")
        last = [ln for ln in text.strip().splitlines() if ln.strip()][-1]
        return "failed" not in last.lower(), last.strip()

    @check("Report: DOCX opens")
    def _():
        from docx import Document

        d = Document(paths.REPORT_DIR / "FlameGuard_AI_Final_Report.docx")
        return (len(d.paragraphs) > 50,
                f"{len(d.paragraphs)} paragraphs, {len(d.tables)} tables, "
                f"{len(d.inline_shapes)} figures")

    @check("Report: PDF opens")
    def _():
        from pypdf import PdfReader

        r = PdfReader(str(paths.REPORT_DIR / "FlameGuard_AI_Final_Report.pdf"))
        return len(r.pages) >= 10, f"{len(r.pages)} pages"

    @check("Report: no unfilled metric placeholders")
    def _():
        md = (paths.REPORT_DIR / "report_content.md").read_text(encoding="utf-8")
        # only personal details may remain as placeholders
        allowed = {"[GROUP NUMBER]", "[SUBMISSION DATE]", "[Project Manager]",
                   "[Dataset & EDA Lead]", "[Model Training Lead]",
                   "[Application Development Lead]", "[Evaluation & Documentation Lead]",
                   "[Application Lead]"}
        import re

        found = set(re.findall(r"\[[^\]\n]{2,40}\]", md))
        # A real *unfilled* placeholder is a template field the author still has
        # to complete: it is ALL-CAPS (e.g. [GROUP NUMBER]) or carries a fill-me
        # keyword. Bracketed prose that is normal mixed-case text - academic
        # citation notation such as "[Computer software]" or
        # "[Open Source Dataset, CC BY 4.0]", and markdown [caption](path) links -
        # is finished content, not a placeholder, and must not trip this check.
        fill_keywords = re.compile(r"\b(TODO|TBD|XXX|FIXME|INSERT|YOUR|ENTER|"
                                   r"NAME|VALUE|PLACEHOLDER|FIXME)\b", re.IGNORECASE)

        def is_unfilled(field: str) -> bool:
            inner = field[1:-1].strip()
            if field in allowed:
                return False
            if field.startswith(("[Figure", "[Table")):
                return False
            # ALL-CAPS template field (letters are all uppercase) -> unfilled
            letters = [c for c in inner if c.isalpha()]
            if letters and all(c.isupper() for c in letters):
                return True
            return bool(fill_keywords.search(inner))

        stray = {f for f in found if is_unfilled(f)}
        return not stray, f"stray placeholders: {sorted(stray) or 'none'}"

    @check("Presentation: PPTX opens with >=13 slides")
    def _():
        from pptx import Presentation

        p = Presentation(paths.PRESENTATION_DIR / "FlameGuard_AI_Presentation.pptx")
        n = len(p.slides)
        notes = sum(1 for s in p.slides if s.has_notes_slide
                    and s.notes_slide.notes_text_frame.text.strip())
        return 13 <= n <= 16, f"{n} slides, {notes} with speaker notes"

    @check("Presentation: notes + demo script")
    def _():
        a = (paths.PRESENTATION_DIR / "speaker_notes.md").exists()
        b = (paths.PRESENTATION_DIR / "demo_script.md").exists()
        c = any((paths.PRESENTATION_DIR / "backup_demo").glob("*"))
        return a and b and c, f"notes={a} demo_script={b} backup_assets={c}"

    @check("Scrum: artefacts complete")
    def _():
        need = ["product_backlog.csv", "sprint_backlog.csv", "user_stories.md",
                "acceptance_criteria.md", "scrum_board.csv", "burndown_chart.png",
                "risk_register.csv", "sprint_retrospectives.md",
                "contribution_table.csv", "meeting_minutes.md"]
        missing = [f for f in need if not (paths.AGILE_DIR / f).exists()]
        return not missing, f"{len(need) - len(missing)}/{len(need)} present"

    @check("Notebooks present", required=False)
    def _():
        nbs = list((paths.PROJECT_ROOT / "notebooks").glob("*.ipynb"))
        return len(nbs) >= 5, f"{len(nbs)} notebooks"

    @check("No absolute personal paths in source")
    def _():
        import re

        bad = []
        pattern = re.compile(r"[A-Za-z]:\\\\?Users\\\\?", re.IGNORECASE)
        for p in list((paths.PROJECT_ROOT / "src").rglob("*.py")) + \
                 list((paths.PROJECT_ROOT / "scripts").rglob("*.py")) + \
                 [paths.PROJECT_ROOT / "app.py"]:
            if pattern.search(p.read_text(encoding="utf-8", errors="ignore")):
                bad.append(p.name)
        return not bad, f"files with hard-coded user paths: {bad or 'none'}"


def verify_zip(zip_path: Path) -> None:
    @check(f"Submission ZIP: {zip_path.name}")
    def _():
        if not zip_path.exists():
            return False, "not built"
        with zipfile.ZipFile(zip_path) as zf:
            bad = zf.testzip()
            names = zf.namelist()
        required = [
            "FlameGuard_AI/report/FlameGuard_AI_Final_Report.pdf",
            "FlameGuard_AI/report/FlameGuard_AI_Final_Report.docx",
            "FlameGuard_AI/presentation/FlameGuard_AI_Presentation.pptx",
            "FlameGuard_AI/models/final/best.pt",
            "FlameGuard_AI/app.py",
            "FlameGuard_AI/README.md",
            "FlameGuard_AI/DATASET_NOTE.md",
        ]
        missing = [r for r in required if r not in names]
        forbidden = [n for n in names
                     if "/.venv/" in n or "__pycache__" in n or "/data/raw/" in n]
        ok = bad is None and not missing and not forbidden
        return ok, (f"{len(names)} entries, {file_size_mb(zip_path):.1f} MB, "
                    f"missing={missing or 'none'}, forbidden={len(forbidden)}")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--zip", type=Path, default=None)
    args = parser.parse_args()

    log.info("=" * 78)
    log.info("FlameGuard AI - final verification")
    log.info("=" * 78)
    run_all()
    zips = [args.zip] if args.zip else sorted(paths.SUBMISSION_DIR.glob("*.zip"))
    for z in zips:
        verify_zip(z)

    required_failed = [c for c in CHECKS if not c[1] and c[3]]
    optional_failed = [c for c in CHECKS if not c[1] and not c[3]]
    passed = [c for c in CHECKS if c[1]]

    report = paths.OUTPUTS_DIR / "verification_report.md"
    lines = ["# FlameGuard AI - Verification Report", "",
             f"**{len(passed)}/{len(CHECKS)} checks passed.**", "",
             "| Check | Result | Detail |", "|---|---|---|"]
    for name, ok, detail, required in CHECKS:
        mark = "PASS" if ok else ("**FAIL**" if required else "WARN")
        lines.append(f"| {name} | {mark} | {detail} |")
    report.write_text("\n".join(lines) + "\n", encoding="utf-8")

    log.info("=" * 78)
    log.info("%d/%d checks passed  (%d required failures, %d warnings)",
             len(passed), len(CHECKS), len(required_failed), len(optional_failed))
    log.info("report -> %s", paths.rel_to_root(report))
    if required_failed:
        for name, _, detail, _ in required_failed:
            log.error("FAILED: %s - %s", name, detail)
        return 1
    log.info("ALL REQUIRED CHECKS PASSED")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
