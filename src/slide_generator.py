"""Presentation deck generation (python-pptx) from measured artefacts.

15 slides for a ~15 minute talk, one message per slide, real charts and real
numbers only.  Speaker notes are attached to each slide and also exported to
presentation/speaker_notes.md.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from src import paths
from src.utils import setup_logging

log = setup_logging("flameguard.slides")

ACCENT = RGBColor(0xB0, 0x3A, 0x2E)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

W, H = Inches(13.333), Inches(7.5)      # 16:9


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title(slide, text: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), W - Inches(1.4), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = DARK
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(15)
        r2.font.color.rgb = GREY
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.55), Inches(1.4), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False


def _bullets(slide, items: list[str], top: float = 2.0, left: float = 0.75,
             width: float | None = None, size: int = 17) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width) if width else W - Inches(1.5),
                                   H - Inches(top + 0.7))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = DARK
        p.space_after = Pt(11)


def _picture(slide, path: Path, left: float, top: float,
             max_w: float, max_h: float) -> None:
    """Insert an image scaled to fit a box, preserving aspect ratio."""
    if not path.exists():
        log.warning("slide image missing: %s", path)
        return
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                             width=Emu(int(iw * scale * 914400)),
                             height=Emu(int(ih * scale * 914400)))


def _metric_cards(slide, cards: list[tuple[str, str]], top: float = 5.35) -> None:
    n = len(cards)
    gap, margin = 0.25, 0.75
    total = 13.333 - 2 * margin
    cw = (total - gap * (n - 1)) / n
    for i, (label, value) in enumerate(cards):
        left = margin + i * (cw + gap)
        shape = slide.shapes.add_shape(5, Inches(left), Inches(top),
                                       Inches(cw), Inches(1.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT
        shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = value
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(11)
        r2.font.color.rgb = GREY


def _notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def build_deck(art: Any) -> tuple[Presentation, list[tuple[str, str, str]]]:
    """Returns the deck and (title, time, notes) rows for speaker_notes.md."""
    tm = art.test_metrics
    fire, smoke = tm["per_class"]["Fire"], tm["per_class"]["Smoke"]
    eda, rs = art.eda, art.resplit
    err = art.errors
    thr = art.model_meta["confidence_threshold"]
    gpu_s = art.speed.get("gpu")
    fast = gpu_s or art.speed["cpu"]
    leak_pct = rs["leakage_before"]["images_in_spanning_groups"] / rs["images"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    notes: list[tuple[str, str, str]] = []

    def add(title, subtitle, note, minutes):
        s = _blank(prs)
        _title(s, title, subtitle)
        _notes(s, note)
        notes.append((title, minutes, note))
        return s

    # 1 - title
    s = _blank(prs)
    box = s.shapes.add_textbox(Inches(1.0), Inches(2.3), W - Inches(2.0), Inches(2.4))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "FlameGuard AI"
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "Real-Time Fire and Smoke Detection Using Transfer Learning"
    r.font.size = Pt(23)
    r.font.color.rgb = DARK
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "AASD 4014 - Deep Learning II  |  Group [GROUP NUMBER]"
    r.font.size = Pt(15)
    r.font.color.rgb = GREY
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "Muthukumaran"
    r.font.size = Pt(12)
    r.font.color.rgb = GREY
    note1 = ("Open: 'Fire kills, and the clock starts at ignition - not at alarm.' "
             "Introduce the team and the one-line goal: detect fire and smoke in "
             "images, video and live camera, using a model we fine-tuned ourselves. "
             "Keep this to 45 seconds.")
    _notes(s, note1)
    notes.append(("Title", "0:45", note1))

    # 2 - problem
    s = add("The problem", "Point sensors wait for smoke to arrive. Cameras do not.",
            "Make the motivation concrete: a smoke detector needs the smoke to physically "
            "reach it, which in a warehouse, an atrium or a forest can take minutes or "
            "never happen. Cameras already watch these spaces. The hard part is that "
            "smoke looks like fog, steam and cloud, and sunsets look like fire - so the "
            "model must balance missing a real fire against crying wolf.", "1:15")
    _bullets(s, [
        "Time-to-alarm is the strongest predictor of fire severity",
        "Point sensors only fire once smoke physically reaches them - slow outdoors and in large spaces",
        "Cameras are already installed nearly everywhere: detection software is the cheap part",
        "But: smoke ≈ fog, steam, cloud.  Fire ≈ sunset, orange lamps, reflections",
        "So the real design question is where to sit on the precision/recall trade-off",
    ])

    # 3 - objective
    s = add("Objective and scope", "What we committed to build",
            "State scope crisply. Two classes, transfer learning (the course forbids "
            "presenting an unchanged pretrained model), a real application, and honest "
            "evaluation. Mention the constraint that shaped everything: one 4GB laptop "
            "GPU.", "1:00")
    _bullets(s, [
        "Detect Fire and Smoke with bounding boxes + confidence, in images, video and live camera",
        "Fine-tune pretrained YOLO weights on a custom dataset - never an unchanged model",
        "Locally-trained model only: no paid inference API",
        "Deliver a working, demonstrable application with downloadable outputs",
        "Constraint that shaped every decision: a single 4GB laptop GPU",
    ])

    # 4 - dataset
    s = add("The dataset", "Roboflow 'fire and smoke' v1 - CC BY 4.0",
            f"Give the shape of the data: {eda['total_images']:,} images, "
            f"{eda['total_annotations']:,} boxes, imbalanced roughly "
            f"{eda['class_imbalance_ratio']:.1f} to 1 toward fire. Emphasise the "
            f"{eda['background_images']:,} background images - they are not waste, they "
            f"are how the model learns not to alarm on a sunset.", "1:00")
    _bullets(s, [
        f"{eda['total_images']:,} images · {eda['total_annotations']:,} annotations · 2 classes (Fire, Smoke)",
        f"Fire in {eda['images_with_fire']:,} images, Smoke in {eda['images_with_smoke']:,}, both in {eda['images_with_both']:,}",
        f"Class imbalance {eda['class_imbalance_ratio']:.2f} : 1 in favour of Fire",
        f"{eda['background_images']:,} background images with no fire or smoke - these teach the model what NOT to flag",
        "Licence CC BY 4.0 · original ZIP preserved unmodified",
    ], top=2.1, width=6.0)
    _picture(s, paths.EDA_OUTPUT_DIR / "02_class_balance.png", 7.0, 2.1, 5.7, 3.6)

    # 5 - leakage (the headline finding)
    s = add("What we found before training anything",
            "The published train/valid/test split leaked",
            f"This is the slide to slow down on. The dataset ships mirrored and "
            f"noise-injected copies of its own images, and it scattered them across "
            f"train, validation and test. {rs['leakage_before']['groups_spanning_splits']} "
            f"groups - {leak_pct:.0%} of all images - had members in more than one split. "
            f"Training on that and reporting test numbers would measure memorisation. We "
            f"grouped images by perceptual hash and rebuilt the splits group-wise. Every "
            f"number in this talk comes from the repaired data. If you take one thing "
            f"from this presentation, take this one.", "2:00")
    _bullets(s, [
        "The export contains its own augmented copies: Mirror… and Noise… duplicates, plus sequential video frames",
        f"Those copies were spread across train / validation / test: {rs['leakage_before']['groups_spanning_splits']} groups, "
        f"{rs['leakage_before']['images_in_spanning_groups']:,} images = {leak_pct:.0%} of the dataset",
        "Evaluating on that split measures memorisation, not generalisation",
        f"Fix: perceptual-hash grouping ({rs['source_groups']:,} source groups) → whole groups assigned to one split, seed {rs['seed']}",
        f"After repair: {rs['leakage_after']['groups_spanning_splits']} groups span splits.  All results below use the repaired data.",
    ])

    # 6 - EDA
    s = add("What the data told us", "EDA drove the modelling choices",
            f"Three findings that changed decisions: {eda['small_boxes_pct']:.0f}% of "
            f"boxes are small, so we kept the full 640px input instead of downscaling "
            f"for speed. Smoke sits high in the frame and fire sits low - which is why we "
            f"disabled vertical flip augmentation: an upside-down plume is not a scene "
            f"that exists. And a tenth of the images are dark night fires, so we kept "
            f"brightness jitter moderate.", "1:15")
    _bullets(s, [
        f"{eda['small_boxes_pct']:.0f}% of objects are small → keep 640px input, enable mosaic",
        "Smoke rises: smoke boxes sit high, fire sits low → vertical flip DISABLED (unphysical)",
        "~10% of images are night scenes → moderate brightness jitter, don't wash out the dark",
        f"Smoke boxes are larger and hazier than fire boxes → harder to localise (AP {smoke['ap50']:.2f} vs {fire['ap50']:.2f})",
    ], top=2.1, width=5.8)
    _picture(s, paths.EDA_OUTPUT_DIR / "04_center_heatmap.png", 6.9, 2.4, 5.8, 3.2)

    # 7 - model
    s = add("The model", "YOLO, fine-tuned - not used as-is",
            "Explain transfer learning in one breath: COCO has 80 classes and none of "
            "them is fire, so the pretrained model cannot detect fire at all. What it "
            "brings is generic vision - edges, textures, shapes - learned from 118,000 "
            "images. We fine-tune ALL layers on our data. Single-stage detection is what "
            "makes the live webcam demo possible.", "1:30")
    _bullets(s, [
        "Single-stage anchor-free detector: boxes + classes in one forward pass → real-time capable",
        "Backbone (CSPDarknet) → Neck (PAN/FPN, fuses fine detail with semantics) → decoupled head at 3 scales",
        "Loss = CIoU box + BCE classification + Distribution Focal Loss;  NMS at inference",
        "Transfer learning: start from COCO weights, fine-tune ALL layers on fire/smoke",
        "COCO has no fire class - the pretrained model detects nothing here until we adapt it",
    ])

    # 8 - experiments
    exp = art.experiments.set_index("experiment_id")
    v8s_done = "e2_stronger_v8s" in exp.index
    if v8s_done:
        note = ("Be honest about the compute wall: FP32 was forced on us because AMP "
                "produces NaN losses on GTX 16xx GPUs. We measured YOLOv8s's memory cost "
                "before training it - batch 8 wants ~7.9GB, batch 4 ~6.1GB, both spilling "
                "to system RAM - so we trained it at batch 2, the only size that fits, "
                "with gradient accumulation to a nominal batch of 64. It completed, and "
                "the bigger backbone still did not beat the baseline at our budget.")
    else:
        note = ("Be honest about the compute wall: FP32 was forced on us because AMP "
                "produces NaN losses on GTX 16xx GPUs, which doubled epoch time. YOLOv8s "
                "needs ~7.9GB at batch 8 - nearly twice our VRAM - so it spilled to system "
                "memory; we report its measured cost rather than a half-trained number.")
    s = add("Experiments", "Baseline, architecture comparison, capacity probe", note, "1:30")
    rows = []
    if "e1_baseline_v8n" in exp.index:
        e1 = exp.loc["e1_baseline_v8n"]
        rows.append(f"YOLOv8n baseline · {int(e1.epochs_run)} epochs · mAP@0.5 = {e1.map50:.3f} · {e1.duration}")
    if "e3_compare_11n" in exp.index:
        e3 = exp.loc["e3_compare_11n"]
        rows.append(f"YOLO11n · same batch/imgsz · {int(e3.epochs_run)} epochs · mAP@0.5 = {e3.map50:.3f}")
    if v8s_done:
        e2 = exp.loc["e2_stronger_v8s"]
        rows.append(f"YOLOv8s (3.3x FLOPs) · {int(e2.epochs_run)} epochs @ batch 2 · mAP@0.5 = {e2.map50:.3f}")
        if art.vram_probe:
            v = {r["batch"]: r for r in art.vram_probe["runs"] if r["model"].startswith("yolov8s")}
            parts = ", ".join(f"b{b}={v[b]['peak_reserved_gb']}GB" for b in sorted(v))
            rows.append(f"…batch 2 was forced by 4GB VRAM (measured: {parts}); bigger backbone ≠ better here")
    elif art.vram_probe:
        v = {r["batch"]: r for r in art.vram_probe["runs"] if r["model"].startswith("yolov8s")}
        parts = " · ".join(f"b{b}: {v[b]['peak_reserved_gb']}GB" for b in sorted(v))
        rows.append(f"YOLOv8s: NOT COMPLETED — measured cost on a 4GB card: {parts}")
        rows.append("Over 4GB it doesn't crash — it pages to system RAM and throughput "
                    "collapses ~5x. Only batch 2 fits.")
    rows += [
        "AMP auto-disabled on GTX 16xx (NaN losses) → FP32 → ~2x slower",
        "Different epoch counts → compared at EQUAL epochs (metrics logged every epoch)",
    ]
    _bullets(s, rows, top=2.05, size=15)

    # 9 - tuning
    s = add("Hyperparameter tuning", "One variable at a time, against a control",
            "The design matters more than the result: a control run at the same budget, "
            "then three probes each changing exactly one thing. That is what makes any "
            "difference attributable. Report the outcome honestly - including probes that "
            "did not beat the control.", "1:15")
    probes = art.experiments[art.experiments.experiment_id.str.startswith("e4")]
    probe_rows = ["Control: default recipe at the probe budget - every comparison is against this"]
    for _, r in probes.iterrows():
        if r.experiment_id == "e4d_probe_baseline":
            continue
        label = {"e4a_probe_adamw": "A · lower learning rate (1.67e-3 → 1e-3)",
                 "e4b_probe_augment": "B · stronger HSV/scale augmentation",
                 "e4c_probe_loss": "C · classification loss weight 0.5 → 1.0"}.get(
            r.experiment_id, r.experiment_id)
        probe_rows.append(f"{label} → mAP@0.5:0.95 = {r.map50_95:.3f}")
    if not probes.empty and "e4d_probe_baseline" in set(probes.experiment_id):
        ctrl = probes[probes.experiment_id == "e4d_probe_baseline"].iloc[0]
        probe_rows.append(f"Control result: mAP@0.5:0.95 = {ctrl.map50_95:.3f}")
    exp_all = art.experiments.drop_duplicates("experiment_id", keep="last").set_index("experiment_id")
    if "e5a_naive_restart" in exp_all.index and "e1_baseline_v8n" in exp_all.index:
        bad = exp_all.loc["e5a_naive_restart"]
        b = exp_all.loc["e1_baseline_v8n"]
        probe_rows.append(
            f"Then the final model FAILED once: restarting a schedule on the converged "
            f"baseline dropped mAP@0.5 {b['map50']:.3f} → {bad['map50']:.3f} "
            f"(best epoch = {int(bad['best_epoch'])})")
        probe_rows.append(
            "Fix: continue, don't restart — low LR, no warm-up, no mosaic")
    probe_rows.append("Caveat: short probes rank settings under a SHORT schedule")
    _bullets(s, probe_rows, top=2.05, size=15)

    # 10 - threshold + final results
    s = add("Final results", "Test split evaluated exactly once",
            f"Stress the discipline: the threshold was chosen on validation, then the "
            f"test set was touched once. Overall mAP@0.5 is {tm['map50']:.3f}. The "
            f"per-class gap is the story - Fire AP {fire['ap50']:.3f} vs Smoke "
            f"{smoke['ap50']:.3f}. Smoke is harder, exactly as the EDA predicted.", "1:45")
    _bullets(s, [
        f"Confidence threshold {thr:.2f} chosen on VALIDATION (F1-optimal), then frozen",
        f"Test split evaluated ONCE, after model and threshold were fixed",
        f"Fire: P {fire['precision']:.3f} · R {fire['recall']:.3f} · AP@0.5 {fire['ap50']:.3f}",
        f"Smoke: P {smoke['precision']:.3f} · R {smoke['recall']:.3f} · AP@0.5 {smoke['ap50']:.3f}",
        f"Speed: {fast['mean_ms']:.0f} ms/image ({fast['fps']:.0f} FPS) on "
        f"{'GPU' if gpu_s else 'CPU'}, measured end-to-end",
    ], top=2.1, size=16)
    _metric_cards(s, [
        ("Precision", f"{tm['precision']:.3f}"),
        ("Recall", f"{tm['recall']:.3f}"),
        ("F1", f"{tm['f1']:.3f}"),
        ("mAP@0.5", f"{tm['map50']:.3f}"),
        ("mAP@0.5:0.95", f"{tm['map50_95']:.3f}"),
    ])

    # 11 - error analysis
    orange_conf = 0.0
    if art.colour_probe:
        orange_conf = art.colour_probe["probes"].get("flat_orange", {}).get(
            "max_confidence", 0.0)
    s = add("Where it fails", "And why that is the useful slide",
            f"Do not skip this. The model makes almost no class confusions "
            f"({tm['confusion_counts']['cross_class_confusions']} in the whole test set) - "
            f"its errors are about whether something is there at all. Then land the "
            f"punchline: we fed it a plain orange rectangle - no fire, no texture, no "
            f"structure - and it said Fire with {orange_conf:.2f} confidence. The model "
            f"has learned a colour prior, not a concept of flame. That one experiment "
            f"explains the entire false-positive gallery, and it is why hard-negative "
            f"mining beats architecture search as the next step.", "1:45")
    _bullets(s, [
        f"{err['total']['tp']:,} true positives · {err['total']['fp']:,} false positives · "
        f"{err['total']['fn']:,} missed · {err['total']['loc']:,} localisation errors",
        f"Only {tm['confusion_counts']['cross_class_confusions']} fire↔smoke confusions in the whole test set "
        f"→ the problem is DETECTION, not CLASSIFICATION",
        f"Diagnostic probe: a plain ORANGE RECTANGLE is detected as Fire at "
        f"{orange_conf:.2f} confidence" if orange_conf else
        "Diagnostic probe: flat colour fields trigger false Fire detections",
        "→ the model learned a COLOUR PRIOR, not the structure of flame",
        "That explains every false positive: sunsets, orange lamps, warm reflections",
        "Fix: hard-negative mining — not a bigger model",
    ], top=2.1, width=6.0, size=15)
    _picture(s, paths.ERROR_ANALYSIS_OUTPUT_DIR / "error_gallery.png", 7.1, 2.1, 5.6, 4.6)

    # 12 - application
    s = add("The application", "One inference engine, three ways in",
            "Point out the engineering: image, video and webcam all call the same predict "
            "path, so a threshold means the same thing everywhere and there is exactly one "
            "place a detection bug could hide. Mention the OpenCV fallback and why it "
            "exists - browser camera access fails in exactly the situation where a demo "
            "must not fail.", "1:15")
    _bullets(s, [
        "Streamlit app · one cached model · shared inference engine for all three modes",
        "Image upload → annotated result + counts + confidences + PNG/CSV/JSON download",
        "Video upload → frame-by-frame with progress, frame-skip control, H.264 output, per-frame CSV",
        "Live camera (streamlit-webrtc) → boxes, live counts, measured FPS, smoothed status banner",
        "OpenCV desktop fallback (python src/webcam_inference.py) for when the browser won't cooperate",
        "Model Performance tab reads saved files - no hard-coded numbers anywhere in the UI",
    ], top=2.1, width=6.2, size=16)
    _picture(s, paths.SCREENSHOTS_DIR / "02_image_detection_result.png", 7.3, 2.3, 5.5, 4.0)

    # 13 - demo
    s = add("Live demonstration", "Image → video → live camera",
            "DEMO ORDER (rehearse this): 1) show model + device status in the sidebar; "
            "2) upload a fire image; 3) upload a smoke image; 4) upload a both-classes "
            "image; 5) upload a negative/difficult image and let it be wrong - explain "
            "why; 6) run the short video; 7) start the webcam; 8) show the performance "
            "tab. If the webcam fails, switch to the OpenCV fallback; if that fails, play "
            "the backup video in presentation/backup_demo/. Never debug live - fall "
            "back.", "2:30")
    _bullets(s, [
        "1 · Model + device status   2 · Fire image   3 · Smoke image   4 · Both classes",
        "5 · A difficult / negative image - and an honest explanation of the failure",
        "6 · Short video with progress + downloads   7 · Live webcam   8 · Performance tab",
        "Backup plan: OpenCV fallback → pre-recorded video → screenshots.  Never debug live.",
    ])

    # 14 - scrum
    s = add("How we worked", "Four sprints, five roles",
            "Keep it short. The point worth making: the risk register earned its keep - "
            "we identified the leakage risk during Sprint 1 planning and mitigated it "
            "before spending GPU time on a model whose evaluation would have been "
            "worthless.", "1:00")
    _bullets(s, [
        "4 one-week sprints · 30 backlog items · 84 story points · 5 defined roles",
        "Sprint 1 data & leakage repair · 2 baseline+comparison · 3 tuning+app · 4 evaluation+delivery",
        "Risk register caught the leakage risk in planning - before we wasted GPU time on it",
        "Artefacts: backlog, board, burndown, retrospectives, contribution table (agile/)",
    ], top=2.1, width=6.3)
    _picture(s, paths.AGILE_DIR / "burndown_chart.png", 7.4, 2.3, 5.4, 3.6)

    # 15 - conclusion
    s = add("Conclusion and future work", "What we learned, what we would do next",
            "Land the three lessons. Then the disclaimer - say it out loud, do not just "
            "leave it on the slide. Finish by inviting questions.", "1:15")
    _bullets(s, [
        f"Built: leakage-free dataset → fine-tuned detector (mAP@0.5 {tm['map50']:.3f}) → working app with live camera",
        "Lesson 1: data hygiene moved our numbers more than any hyperparameter did",
        "Lesson 2: the model's errors are detection errors, not classification errors - that tells us what to fix",
        "Lesson 3: the confidence threshold is a safety decision, not a library default",
        "Next: thermal/IR fusion · temporal modelling (smoke moves, clouds don't) · hard-negative mining · edge deployment",
        "FlameGuard AI is an educational prototype - NOT a certified fire-detection system",
    ], top=2.0, size=16)
    return prs, notes


def write_speaker_notes(notes: list[tuple[str, str, str]], out: Path) -> None:
    total = sum(int(m.split(":")[0]) * 60 + int(m.split(":")[1]) for _, m, _ in notes)
    lines = ["# FlameGuard AI - Speaker Notes", "",
             f"Target length: ~15 minutes (planned {total // 60}m {total % 60}s "
             f"across {len(notes)} slides, leaving room for questions).", ""]
    for i, (title, minutes, note) in enumerate(notes, start=1):
        lines += [f"## Slide {i} - {title}  ({minutes})", "", note, ""]
    lines += ["## If something goes wrong", "",
              "- Webcam will not start → switch to `python src/webcam_inference.py`.",
              "- That also fails → play `presentation/backup_demo/` video and screenshots.",
              "- Video processing is slow → set frame skip to 'every 3rd frame' first.",
              "- Never debug live. Narrate the fallback and keep moving.", ""]
    out.write_text("\n".join(lines), encoding="utf-8")
    log.info("speaker notes -> %s", out)


def write_demo_script(art: Any, out: Path) -> None:
    thr = art.model_meta["confidence_threshold"]
    text = f"""# FlameGuard AI - Classroom Demo Script

**Before the session (do this, do not skip it)**

1. `cd FlameGuard_AI` and activate the environment.
2. `python -m pytest tests/ -q` - confirm the suite passes.
3. `streamlit run app.py` - confirm the app loads and the sidebar shows
   "Model loaded" and the correct device.
4. Have `outputs/sample_inputs/` open in a file browser - the demo images are there.
5. Have `presentation/backup_demo/` ready in case the live demo fails.
6. Plug in the laptop. GPU inference on battery is throttled.

**Demo sequence (about 4 minutes)**

| # | Action | What to say |
|---|--------|-------------|
| 1 | Show the sidebar | "The model is loaded once and cached; it is running on the GPU. Confidence is at {thr:.2f} - we chose that on validation data, not by accepting a default." |
| 2 | Image tab → upload a fire-only sample | "Original on the left, detections on the right. Counts, peak confidence and inference time are measured, not estimated." |
| 3 | Upload a smoke-only sample | "Smoke is the harder class - it is semi-transparent and has no crisp boundary. Watch the box hug the dense core." |
| 4 | Upload a both-classes sample | "Fire and smoke together, which is the realistic case: smoke is usually the earlier signal." |
| 5 | Upload a negative or difficult sample | "This one it gets wrong - and this is the slide I want you to remember. It reads the cloud bank as smoke. We know this failure mode because we measured it." |
| 6 | Drag the confidence slider up | "Raising the threshold kills that false positive - and also starts killing real detections. That trade-off is the actual product decision." |
| 7 | Video tab → short clip, frame skip = 2 | "One frame at a time, so memory stays flat regardless of video length. Progress bar is real. The per-frame CSV downloads with timestamps." |
| 8 | Live Camera → START, allow permission | "This is the same engine, running on the webcam feed. Fire and smoke counts and measured FPS update live. The status banner is smoothed over five frames so it does not strobe." |
| 9 | Hold up a phone showing a fire image | "It fires on a screen showing fire - which is worth noting: the model detects the visual signature of fire, not fire itself." |
| 10 | Model Performance tab | "Every number here is read from the evaluation files. If we had not run the pipeline, it would say so rather than showing a number." |
| 11 | About tab | Read the disclaimer aloud. |

**Fallbacks (rehearse these too)**

- Browser will not grant camera access → `python src/webcam_inference.py`
  (native OpenCV window, press Q to quit).
- No camera at all → play the pre-processed video from `presentation/backup_demo/`.
- Everything is slow → set frame skip to "every 3rd frame", drop to a smaller clip.
- The app will not start → walk through `outputs/sample_outputs/sample_grid.png`
  and the screenshots in `outputs/application_screenshots/`.

**Never** debug live in front of the class. Announce the fallback and continue.
"""
    out.write_text(text, encoding="utf-8")
    log.info("demo script -> %s", out)
